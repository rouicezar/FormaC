from pathlib import Path

from pptx import Presentation

from app.ingestion.parsers.base import ParsedSection


class PptxParser:
    def parse(self, path: Path) -> list[ParsedSection]:
        presentation = Presentation(path)
        sections: list[ParsedSection] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                sections.append(
                    ParsedSection(
                        text="\n".join(texts),
                        locator={"slide": slide_number},
                    )
                )
        return sections

