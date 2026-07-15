from collections.abc import Callable
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen.canvas import Canvas

from app.ingestion.chunking import chunk_sections
from app.ingestion.parsers.base import ParsedSection
from app.ingestion.parsers.registry import parse_document


def write_pdf(path: Path) -> None:
    canvas = Canvas(str(path))
    canvas.drawString(72, 720, "PDF support anchor")
    canvas.showPage()
    canvas.save()


def write_docx(path: Path) -> None:
    document = Document()
    document.add_paragraph("DOCX policy anchor")
    document.save(path)


def write_markdown(path: Path) -> None:
    path.write_text("# Guide\n\nMarkdown guide anchor", encoding="utf-8")


def write_text(path: Path) -> None:
    path.write_text("Plain text answer anchor", encoding="utf-8")


def write_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "FAQ"
    sheet.append(["Question", "Answer"])
    sheet.append(["Shipping", "XLSX shipping anchor"])
    workbook.save(path)


def write_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    text_box = slide.shapes.add_textbox(0, 0, 4_000_000, 1_000_000)
    text_box.text = "PPTX onboarding anchor"
    presentation.save(path)


FixtureWriter = Callable[[Path], None]


@pytest.mark.parametrize(
    ("suffix", "writer", "anchor", "locator_key", "locator_value"),
    [
        (".pdf", write_pdf, "PDF support anchor", "page", 1),
        (".docx", write_docx, "DOCX policy anchor", "paragraph", 1),
        (".md", write_markdown, "Markdown guide anchor", "line_start", 1),
        (".txt", write_text, "Plain text answer anchor", "line_start", 1),
        (".xlsx", write_xlsx, "XLSX shipping anchor", "sheet", "FAQ"),
        (".pptx", write_pptx, "PPTX onboarding anchor", "slide", 1),
    ],
)
def test_supported_document_parser_preserves_location(
    tmp_path: Path,
    suffix: str,
    writer: FixtureWriter,
    anchor: str,
    locator_key: str,
    locator_value: str | int,
) -> None:
    path = tmp_path / f"fixture{suffix}"
    writer(path)

    sections = parse_document(path)

    assert sections
    matching = [section for section in sections if anchor in section.text]
    assert matching
    assert matching[0].locator[locator_key] == locator_value


def test_chunking_normalizes_whitespace_and_preserves_locator() -> None:
    sections = [
        ParsedSection(
            text="First   answer\n\nSecond answer",
            locator={"page": 3},
        )
    ]

    chunks = chunk_sections(sections, max_characters=20)

    assert [chunk.chunk_index for chunk in chunks] == [0, 1]
    assert chunks[0].text == "First answer"
    assert chunks[1].text == "Second answer"
    assert all(chunk.locator == {"page": 3} for chunk in chunks)
